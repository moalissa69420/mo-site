#!/usr/bin/env python3
"""
Build pipeline for Google — Image Equity case study.

- Parses PPTX in ./Google/*.pptx via python-pptx
- Exports slide PNGs (white background) to ./Google/exports/
- Extracts slide texts (simple bulletization)
- Builds manifest.json grouped by section ids matching our case study structure

Note: This is a best-effort simple parser; it avoids OCR and PDF logic.
"""
import json
import os
import sys
from glob import glob
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as e:
    print("python-pptx is required: pip install python-pptx", file=sys.stderr)
    raise

try:
    from PIL import Image
except Exception:
    Image = None  # Exporting thumbnails only if Pillow is available


ROOT = Path(__file__).resolve().parent
GOOGLE_DIR = ROOT / "Google"
EXPORTS = GOOGLE_DIR / "exports"
PICTURES = GOOGLE_DIR / "pictures"
MANIFEST = GOOGLE_DIR / "manifest.json"


SECTION_IDS = [
    "info","exec","team","market","consumer","problem","solution","visual","competitive",
    "location","model","financial","unit","funding","scenarios","revenue","fiveyear","license",
    "menu","brand","ops","timeline","risks","expansion","growth","highlights","ask","appendix"
]


def extract_texts_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        if not hasattr(shape, "has_text_frame"):
            continue
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            s = "".join(run.text for run in p.runs).strip()
            if s:
                texts.append(s)
    # Deduplicate, keep order
    seen = set()
    out = []
    for t in texts:
        if t not in seen:
            seen.add(t)
            out.append(t)
    return out[:20]


def export_slide_image(prs, slide_index, out_dir):
    """
    python-pptx does not render slides; we create a blank placeholder PNG per slide.
    For now, we return a placeholder path so the UI can render structure.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    filename = out_dir / f"slide_{slide_index+1:02d}.png"
    if Image and not filename.exists():
        img = Image.new("RGB", (1600, 900), color=(255, 255, 255))
        img.save(filename)
    # Return a web path that is resolvable from site root
    return f"Google/exports/{filename.name}"


def extract_pictures_from_slide(slide, slide_index):
    paths = []
    PICTURES.mkdir(parents=True, exist_ok=True)
    pic_idx = 0
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext or 'png'
                pic_idx += 1
                filename = PICTURES / f"pic_{slide_index+1:02d}_{pic_idx}.{ext}"
                if not filename.exists():
                    with open(filename, 'wb') as f:
                        f.write(image.blob)
                paths.append(f"Google/pictures/{filename.name}")
        except Exception:
            # Skip any problematic shapes silently
            continue
    return paths


def infer_section_id_by_index(idx):
    # Simple sequential mapping: first N slides fall into sections in order
    return SECTION_IDS[min(idx, len(SECTION_IDS)-1)]


def extract_toc_titles(prs):
    """Attempt to extract Table of Contents/Agenda slide item titles.
    Returns a list of strings, in order, up to the number of sections.
    """
    keywords = {"table of contents", "contents", "agenda", "overview"}
    def all_text(slide):
        lines = []
        for sh in slide.shapes:
            if hasattr(sh, "has_text_frame") and sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    s = "".join(run.text for run in p.runs).strip()
                    if s:
                        lines.append(s)
        return lines
    for slide in prs.slides:
        lines = all_text(slide)
        joined_lower = " ".join(lines).lower()
        if any(k in joined_lower for k in keywords) and len(lines) >= 3:
            # Heuristic: drop the first line if it looks like the header
            candidates = lines[1:]
            # Keep concise items only
            cleaned = []
            seen = set()
            for t in candidates:
                t = t.strip().strip("-•· ")
                if not t or t.lower() in keywords:
                    continue
                if t in seen:
                    continue
                seen.add(t)
                cleaned.append(t)
            if cleaned:
                return cleaned[:len(SECTION_IDS)]
    return []


def build_manifest(pptx_path: Path):
    prs = Presentation(str(pptx_path))
    sections = { sid: {"title": "", "texts": [], "slides": []} for sid in SECTION_IDS }
    toc_titles = extract_toc_titles(prs)
    for i, slide in enumerate(prs.slides):
        sid = infer_section_id_by_index(i)
        texts = extract_texts_from_slide(slide)
        # Prefer embedded pictures; fallback to blank export
        pic_paths = extract_pictures_from_slide(slide, i)
        if pic_paths:
            images = pic_paths
        else:
            images = [export_slide_image(prs, i, EXPORTS)]
        entry = {
            "slide": i+1,
            "title": (getattr(slide.shapes, 'title', None).text.strip() if getattr(slide.shapes, 'title', None) and getattr(slide.shapes.title, 'text', '') else (texts[0] if texts else f"Slide {i+1}")),
            "images": images
        }
        sections[sid]["slides"].append(entry)
        # add remaining texts as bullets
        if texts:
            # Use up to first 6 relevant lines per slide to keep page readable
            sections[sid]["texts"].extend(texts[:6])
        # Set section display title if not already set
        if not sections[sid]["title"] and entry["title"]:
            sections[sid]["title"] = entry["title"]

    # Apply TOC titles to index labels if available
    if toc_titles:
        for idx, t in enumerate(toc_titles):
            if idx < len(SECTION_IDS):
                sections[SECTION_IDS[idx]]["title"] = t

    # If this is the provided MODU deck, override titles and bullets per spec
    name_lower = pptx_path.name.lower()
    if "modu-google-ppx" in name_lower:
        overrides = [
            {
                "title": "The Infinite Palette",
                "texts": [
                    "An immersive experience combining collaboration and individualism",
                    "Aims to set a new standard for image equity",
                    "Prepared by Abdou Sarr & Mo Alissa — Dec 10, 2021",
                ],
            },
            {
                "title": "The Shirley Card",
                "texts": [
                    "Photography as subjective technology, not just light calibration",
                    "Original Shirley card set color correction standard, excluding diverse skin tones",
                ],
            },
            {
                "title": "Background — The evolution of standard",
                "texts": [
                    "Since the 1990s, major strides improved how cameras see color",
                    "New standards emerged, but no single standard fits infinite skin tones",
                ],
            },
            {
                "title": "Introduction",
                "texts": [
                    "We are more than a color code",
                    "Goal: highlight that the new standard is to not conform to a standard",
                    "Attendees form a human palette in real time by contributing their tones and names",
                ],
            },
            {
                "title": "The Capture Experience",
                "texts": [
                    "Large color grid canvas + booth with Pixel app",
                    "Attendees capture a selfie, add their name, submit to the digital piece",
                    "Optional: show legacy Shirley-card corrected preview",
                ],
            },
            {
                "title": "The AR Experience — Viewing the Digital Palette",
                "texts": [
                    "Canvas IRL is a color grid (old standard); AR app brings it to life",
                    "As submissions accrue, the infinite palette forms in real time",
                    "Viewers see a grid of people, skin tones as backgrounds, names below",
                ],
            },
            {
                "title": "The Purpose — Our Objective",
                "texts": [
                    "Educate: past & present via non-digital experience",
                    "Absorb: capture & submission app builds personal connection",
                    "Immerse: interactive AR evolves as more content arrives",
                    "Embody: leave with a new perspective on color identity",
                ],
            },
            {
                "title": "About MODU",
                "texts": [
                    "Mo Alissa — creative direction; Canon, The Weeknd, Drake",
                    "Abdou Sarr — product management; Shopify, Motorola ATAP",
                    "Work for Nike, Apple, Converse, Australian Open, Amazon, OVO",
                ],
            },
            {
                "title": "Venue — DTLA · SoLa Beehive",
                "texts": [
                    "~120 Student/Org Partner guests (Morning)",
                    "~100 VIP guests (Afternoon)",
                ],
            },
            {
                "title": "User Flow",
                "texts": [
                    "1 Valet Drop‑off · 2 Check‑in · 3 Seen in Power · 4 Seen in Color",
                    "5 Seen in Focus · 6 Seen as You See Yourself · 7 Gifting",
                    "8 Restrooms · 9 Valet Pick‑up",
                ],
            },
            { "title": "Images from Event", "texts": ["Event highlights"] },
            { "title": "Images from Event", "texts": ["Event highlights"] },
        ]
        for idx, ov in enumerate(overrides):
            if idx >= len(SECTION_IDS):
                break
            sid = SECTION_IDS[idx]
            if ov.get("title"):
                sections[sid]["title"] = ov["title"]
            if ov.get("texts"):
                sections[sid]["texts"] = ov["texts"]

    # Trim bullet explosion per section
    for sid in sections:
        uniq = []
        seen = set()
        for t in sections[sid]["texts"]:
            if t in seen:
                continue
            seen.add(t)
            uniq.append(t)
        sections[sid]["texts"] = uniq[:12]

    manifest = {
        "source": pptx_path.name,
        "sections": sections
    }
    MANIFEST.write_text(json.dumps(manifest, indent=2))
    print(f"Wrote {MANIFEST}")


def main():
    # Prefer the new modu deck if present, else support Amazon deck export task
    preferred = GOOGLE_DIR / "modu-google-ppx.pptx"
    if preferred.exists():
        candidates = [str(preferred)]
    else:
        candidates = glob(str(GOOGLE_DIR / "*.pptx"))
    if not candidates:
        print("No PPTX found in ./Google", file=sys.stderr)
        sys.exit(1)
    pptx_path = Path(candidates[0])
    build_manifest(pptx_path)


if __name__ == "__main__":
    main()


